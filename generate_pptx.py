from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_flightguard_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    DARK_BLUE = RGBColor(10, 25, 47)
    ACCENT_BLUE = RGBColor(59, 130, 246)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(226, 232, 240)
    
    def add_title_slide(title, subtitle=""):
        """Slide 1: Cover"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        shapes = slide.shapes
        
        # Background
        background = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = DARK_BLUE
        background.line.fill.background()
        
        # Title
        title_box = shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = WHITE
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(20)
            subtitle_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Team info
        team_box = shapes.add_textbox(Inches(2), Inches(6), Inches(6), Inches(0.6))
        team_frame = team_box.text_frame
        team_frame.text = "Prepared by: Elturan Aliyev | MVP Stage: Testnet Ready"
        team_frame.paragraphs[0].font.size = Pt(14)
        team_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
        team_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, bullets):
        """Generic content slide with bullets"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        
        # Add bullets
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_BLUE
        
        return slide
    
    def add_table_slide(title, headers, rows):
        """Slide with table"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        
        # Title
        title_shape = shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Add table
        table = shapes.add_table(len(rows) + 1, len(headers), Inches(1), Inches(2), Inches(8), Inches(4)).table
        
        # Headers
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
        
        # Data rows
        for r_idx, row in enumerate(rows):
            for c_idx, cell_data in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(cell_data)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        return slide
    
    # SLIDE 1: Cover
    add_title_slide(
        "FlightGuard",
        "Blockchain-Powered Flight Insurance That Pays You Back Automatically"
    )
    
    # SLIDE 2: Problem
    add_content_slide(
        "The Problem: Traditional Insurance Fails Travelers",
        [
            "Energy Waste: Pure PoW blockchains consume massive energy",
            "Centralization Risk: PoS systems concentrate power in few validators",
            "Claim Delays: Traditional insurance takes days/weeks for payouts",
            "Trust Issues: Opaque settlement processes",
            "High Fees: Manual processing drives up costs",
            "Methodology: SCAMPER + 5 Whys analysis identified root causes"
        ]
    )
    
    # SLIDE 3: Solution
    add_content_slide(
        "Our Solution: Hybrid PoW + PoS Consensus",
        [
            "Security Fusion: PoW immutability + PoS speed",
            "Scalability: Next.js SSR + FastAPI async I/O",
            "Inclusivity: Real-time ETL pipelines for instant data",
            "Interoperability: GraphQL + Web3 integration",
            "Meta-Benefit: Trust by Design via methodological rigor (TRIZ, FMEA)"
        ]
    )
    
    # SLIDE 4: Market Analysis
    add_table_slide(
        "Market Segmentation & Use Cases",
        ["Segment", "Use Case", "TAM"],
        [
            ["B2B", "Enterprise supply chain integrations", "$200B"],
            ["B2C", "NFT ticketing + staking rewards", "$150B"],
            ["B2G", "Public records with compliance", "$80B"],
            ["C2C", "P2P exchanges via WebSocket", "$100B"],
            ["G2G", "Cross-border data with oracles", "$50B"]
        ]
    )
    
    # SLIDE 5: Market Opportunity & Economics
    add_content_slide(
        "Market Opportunity & Key Metrics",
        [
            "TAM: $1.5T global blockchain insurance market",
            "SAM: $200M initial capture (2026-2027)",
            "ROI: 15% projected return for stakers",
            "CLV: $500 average Customer Lifetime Value over 3 years",
            "CAC: $50 Customer Acquisition Cost",
            "Why Now: Post-2024 regulatory clarity + DeFi adoption surge"
        ]
    )
    
    # SLIDE 6: Product / Technology
    add_content_slide(
        "Technology Architecture",
        [
            "Frontend: Next.js 16 + React 19 + Three.js (SSR for SEO)",
            "Backend: Hybrid Django 5.0 + FastAPI (Async I/O)",
            "Smart Contracts: Solidity 0.8.20 (Hardhat tested, 100% pass rate)",
            "ETL: Event listeners bridging blockchain to PostgreSQL",
            "TRL: Level 6 (Testnet validated)",
            "Security: FMEA + ReentrancyGuard (RPN reduced 90→10)"
        ]
    )
    
    # SLIDE 7: Business Model
    add_content_slide(
        "Business Model & Tokenomics",
        [
            "Revenue Streams: 0.1% transaction fees, 5% staking commissions",
            "Token Distribution: 20% Team, 30% Community, 25% Treasury",
            "Supply: 1B total, 100M initial circulation",
            "Inflation: 3% annual, halving every 2 years",
            "Utility: Governance voting, staking yields, fee payments",
            "Deflation: Transaction fee burns"
        ]
    )
    
    # SLIDE 8: Business Plan
    add_content_slide(
        "Compliance & Scalability Roadmap",
        [
            "KYC/AML: Optional zero-knowledge proofs for privacy",
            "Regulatory: GDPR, FATF, MiCA compliant",
            "Audits: CertiK + Quantstamp smart contract reviews",
            "Scalability: Cross-chain bridges (IBC-style relay)",
            "Risk Mitigation: FMEA prioritization (RPN <150)",
            "Go-Live: Q1 2026 testnet → Q4 2026 mainnet"
        ]
    )
    
    # SLIDE 9: Traction
    add_content_slide(
        "MVP Traction & Validation",
        [
            "Testnet Deployment: 500+ validators onboarded",
            "Uptime: >99.5% network availability",
            "Gas Efficiency: Average 78,746 gas per transaction",
            "Test Coverage: 100% pass rate (15 passing tests)",
            "Transaction Processing: 2,000+ TPS capacity",
            "Community: 10K+ Discord members, 5K GitHub stars"
        ]
    )
    
    # SLIDE 10: Milestone Plan
    add_table_slide(
        "Roadmap & KPIs (2026-2027)",
        ["Quarter", "Milestone", "KPI"],
        [
            ["Q1 2026", "Testnet MVP Launch", "500 validators, 99.5% uptime"],
            ["Q2 2026", "Third-Party Audits Complete", "<2% issue severity"],
            ["Q3 2026", "Mainnet Beta + DAO Integration", "2,000 TPS, 10K voters"],
            ["Q4 2026", "Treasury Launch", "100 validators, $10M volume"],
            ["Q1 2027", "Institutional Onboarding", "3 pilot governments (B2G)"],
            ["Q3 2027", "AI Security Suite Live", ">95% threat detection"]
        ]
    )
    
    # SLIDE 11: Go-To-Market Strategy
    add_content_slide(
        "Go-To-Market: Phased Rollout",
        [
            "Phase 1: Developer Grants (Q1 2026) - $500K in bounties",
            "Phase 2: Enterprise Integrations (Q2 2026) - B2B/B2G pilots",
            "Phase 3: Cross-Chain Partnerships (Q3 2026) - Polkadot/Cosmos",
            "Phase 4: Retail Yield Portals (Q4 2026) - B2C/C2C campaigns",
            "Marketing Channels: Hackathons, Dev bounties, Institutional reports",
            "Community Building: Airdrops via Wallet Connect, SDK grants"
        ]
    )
    
    # SLIDE 12: Competition & Differentiation
    add_table_slide(
        "Competitive Edge",
        ["Feature", "FlightGuard", "Ethereum", "Polkadot"],
        [
            ["Consensus", "Hybrid PoW+PoS", "PoS Only", "PoS Only"],
            ["Energy Savings", "40% vs PoW", "Baseline", "Similar"],
            ["Security", "Dual-layer", "Single-layer", "Shared"],
            ["AI Integration", "Predictive models", "None", "None"],
            ["Gas Costs", "$0.05 avg", "$2-5", "$0.10"],
            ["Governance", "DAO + Slashing", "DAO", "Council"]
        ]
    )
    
    # SLIDE 13: Team
    add_content_slide(
        "Team & Expertise",
        [
            "Core Team: 5 blockchain engineers, 4 AI scientists, 2 economists",
            "Advisory Board: Stanford blockchain faculty + FinTech veterans",
            "Experience: Ex-Ethereum Foundation, Google AI, McKinsey",
            "Due Care Charter: Ethical disclosure + quarterly audits",
            "Developer Community: 50+ open-source contributors",
            "Legal Advisors: 2 regulatory compliance experts"
        ]
    )
    
    # SLIDE 14: Financials
    add_content_slide(
        "Financial Projections",
        [
            "2026: $10M revenue, 50K users (Testnet + early staking)",
            "2027: $45M revenue, 250K users (DAO + cross-chain)",
            "2028: $100M revenue, 1M users (Institutional integration)",
            "Funding: $5M seed closed, $10M Series A (Q2 2026)",
            "Runway: 24-28 months (extendable via treasury staking)",
            "Unit Economics: CAC $50, CLV $500 (10x return)"
        ]
    )
    
    # SLIDE 15: Due Diligence & Due Care
    add_table_slide(
        "Due Diligence & Due Care Framework",
        ["Category", "Due Diligence (Pre-Launch)", "Due Care (Ongoing)"],
        [
            ["Legal", "Token classification review", "Quarterly compliance audits"],
            ["Technical", "Penetration testing", "Continuous fuzzing (Echidna)"],
            ["Financial", "Treasury audits", "Real-time monitoring dashboards"],
            ["Operational", "Vendor risk checks", "SIEM/SOAR automated responses"],
            ["Governance", "DAO bylaws drafted", "Monthly governance reviews"]
        ]
    )
    
    # SLIDE 16: The Ask
    add_content_slide(
        "The Ask: $10M Series A",
        [
            "Allocation: 40% Development, 30% Marketing, 20% Operations, 10% Reserve",
            "Milestones: Q3 2026 Mainnet, Q4 2026 DAO Treasury",
            "Valuation: $50M pre-money (10x upside potential)",
            "Use of Funds: Smart contract audits, enterprise BD, global expansion",
            "Exit Strategy: Strategic acquisition or token public listing (2028)",
            "Investor Benefits: Early governance rights + staking yields"
        ]
    )
    
    # SLIDE 17: Vision / FOMO
    add_title_slide(
        "Don't Miss the Hybrid Revolution",
        "The first blockchain ecosystem built for intelligence, sustainability, and global interoperability. Missed Bitcoin in 2010? Don't miss this."
    )
    
    prs.save('FlightGuard_Pitch_Deck.pptx')
    print("FlightGuard_Pitch_Deck.pptx created successfully!")

if __name__ == "__main__":
    create_flightguard_presentation()
